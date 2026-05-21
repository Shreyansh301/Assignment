import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# -----------------------------------------------------------
# Load data
# -----------------------------------------------------------
df = pd.read_csv("Hotel_bookings_final.csv")

df["booking_date"]  = pd.to_datetime(df["booking_date"],  errors="coerce")
df["check_in_date"] = pd.to_datetime(df["check_in_date"], errors="coerce")
df["check_out_date"]= pd.to_datetime(df["check_out_date"],errors="coerce")

df["stay_length"]    = (df["check_out_date"] - df["check_in_date"]).dt.days
df["booking_month"]  = df["booking_date"].dt.month

confirmed = df[df["booking_status"] == "Confirmed"]
cancelled = df[df["booking_status"] == "Cancelled"]

# -----------------------------------------------------------
# Quick summary
# -----------------------------------------------------------
print("Total bookings:", len(df))
print("Confirmed:", len(confirmed), f"({len(confirmed)/len(df)*100:.1f}%)")
print("Cancelled:", len(cancelled), f"({len(cancelled)/len(df)*100:.1f}%)")
print("Total revenue: $", round(confirmed["booking_value"].sum() / 1e6, 1), "M")
print("Avg stay length:", round(df["stay_length"].mean(), 1), "days")

# cancellation rate helper
def cancel_rate(group_col):
    return df.groupby(group_col)["booking_status"].apply(
        lambda x: round((x == "Cancelled").sum() / len(x) * 100, 1)
    )

print("\nCancellation by channel:\n",  cancel_rate("booking_channel"))
print("\nCancellation by room type:\n", cancel_rate("room_type"))
print("\nCancellation by star rating:\n",cancel_rate("star_rating"))

# -----------------------------------------------------------
# Chart 1 — Cancellation rates by channel and room type
# -----------------------------------------------------------
fig, axes = plt.subplots(1, 3, figsize=(14, 4))
fig.suptitle("Booking Overview", fontsize=13, fontweight="bold")

# booking status pie
status_counts = df["booking_status"].value_counts()
axes[0].pie(status_counts, labels=status_counts.index, autopct="%1.1f%%",
            colors=["#2ecc71", "#e74c3c", "#95a5a6"], startangle=90)
axes[0].set_title("Booking Status")

# cancellation by channel
chan = cancel_rate("booking_channel")
axes[1].bar(chan.index, chan.values, color=["#3498db", "#e67e22", "#2ecc71"])
axes[1].set_title("Cancellation Rate by Channel (%)")
axes[1].set_ylabel("%")
for i, v in enumerate(chan.values):
    axes[1].text(i, v + 0.3, f"{v}%", ha="center", fontsize=9)

# cancellation by room type
room = cancel_rate("room_type")
axes[2].bar(room.index, room.values, color=["#9b59b6", "#3498db", "#e74c3c"])
axes[2].set_title("Cancellation Rate by Room Type (%)")
axes[2].set_ylabel("%")
for i, v in enumerate(room.values):
    axes[2].text(i, v + 0.3, f"{v}%", ha="center", fontsize=9)

plt.tight_layout()
plt.savefig("chart1_overview.png", dpi=130, bbox_inches="tight")
plt.close()
print("\nSaved chart1_overview.png")

# -----------------------------------------------------------
# Chart 2 — Monthly trends
# -----------------------------------------------------------
monthly_bookings = df.groupby("booking_month").size()
monthly_cancel   = df.groupby("booking_month")["booking_status"].apply(
    lambda x: (x == "Cancelled").sum() / len(x) * 100
)
monthly_revenue  = confirmed.groupby("booking_month")["booking_value"].sum() / 1e6

month_names = ["Jan","Feb","Mar","Apr","May","Jun",
               "Jul","Aug","Sep","Oct","Nov","Dec"]

fig, axes = plt.subplots(1, 3, figsize=(14, 4))
fig.suptitle("Monthly Trends", fontsize=13, fontweight="bold")

axes[0].bar(range(1, 13), monthly_bookings.values, color="#3498db")
axes[0].set_title("Bookings per Month")
axes[0].set_xticks(range(1, 13))
axes[0].set_xticklabels(month_names, rotation=45, fontsize=8)
axes[0].set_ylabel("Number of Bookings")

axes[1].plot(range(1, 13), monthly_cancel.values, marker="o", color="#e74c3c")
axes[1].set_title("Cancellation Rate by Month (%)")
axes[1].set_xticks(range(1, 13))
axes[1].set_xticklabels(month_names, rotation=45, fontsize=8)
axes[1].set_ylabel("%")

axes[2].bar(range(1, 13), monthly_revenue.values, color="#2ecc71")
axes[2].set_title("Monthly Revenue — Confirmed ($M)")
axes[2].set_xticks(range(1, 13))
axes[2].set_xticklabels(month_names, rotation=45, fontsize=8)
axes[2].set_ylabel("Revenue ($M)")

plt.tight_layout()
plt.savefig("chart2_trends.png", dpi=130, bbox_inches="tight")
plt.close()
print("Saved chart2_trends.png")

# -----------------------------------------------------------
# Chart 3 — Revenue and profitability
# -----------------------------------------------------------
rev_by_channel = confirmed.groupby("booking_channel")["booking_value"].sum() / 1e6
val_by_star    = confirmed.groupby("star_rating")["booking_value"].mean()
coupon_cancel  = df.groupby("Coupon USed?")["booking_status"].apply(
    lambda x: round((x == "Cancelled").sum() / len(x) * 100, 1)
)

fig, axes = plt.subplots(1, 3, figsize=(14, 4))
fig.suptitle("Revenue & Profitability", fontsize=13, fontweight="bold")

axes[0].barh(rev_by_channel.index, rev_by_channel.values,
             color=["#3498db", "#e67e22", "#2ecc71"])
axes[0].set_title("Revenue by Channel ($M)")
axes[0].set_xlabel("Revenue ($M)")
for i, v in enumerate(rev_by_channel.values):
    axes[0].text(v + 1, i, f"${v:.0f}M", va="center", fontsize=9)

axes[1].bar([f"{s}★" for s in val_by_star.index], val_by_star.values,
            color=["#95a5a6", "#3498db", "#2ecc71", "#e67e22"])
axes[1].set_title("Avg Booking Value by Star Rating ($)")
axes[1].set_ylabel("Avg Booking Value ($)")

axes[2].bar(["No Coupon", "Coupon Used"], coupon_cancel.values,
            color=["#3498db", "#e67e22"])
axes[2].set_title("Coupon Use vs Cancellation Rate (%)")
axes[2].set_ylabel("%")
for i, v in enumerate(coupon_cancel.values):
    axes[2].text(i, v + 0.2, f"{v}%", ha="center", fontsize=10)

plt.tight_layout()
plt.savefig("chart3_revenue.png", dpi=130, bbox_inches="tight")
plt.close()
print("Saved chart3_revenue.png")

# -----------------------------------------------------------
# Build PowerPoint
# -----------------------------------------------------------
def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_text_box(slide, text, left, top, width, height,
                 size=12, bold=False, color=(30, 30, 30), wrap=True):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(*color)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
blank = prs.slide_layouts[6]

# -- Slide 1: Title --
s = prs.slides.add_slide(blank)
bg = s.background; bg.fill.solid()
bg.fill.fore_color.rgb = rgb(13, 79, 140)

add_text_box(s, "Hotel Bookings — Data Analysis",
             0.5, 1.5, 9, 1.2, size=32, bold=True, color=(255,255,255))
add_text_box(s, "Booking Trends  •  Cancellation Patterns  •  Recommendations",
             0.5, 2.9, 9, 0.6, size=12, color=(180, 210, 240))

stats = [f"30,000 bookings   |   $550.8M revenue   |   20.2% cancellation rate   |   4.0 days avg stay"]
add_text_box(s, stats[0], 0.5, 4.2, 9, 0.5, size=10, color=(180,210,240))

# -- Slide 2: Key Findings --
s = prs.slides.add_slide(blank)
bg = s.background
bg.fill.solid()
bg.fill.fore_color.rgb = rgb(244, 247, 252)

add_text_box(
    s,
    "Key Observations",
    0.4,
    0.2,
    9,
    0.6,
    size=20,
    bold=True,
    color=(13,79,140)
)

findings = [
    "1. Web bookings generated the highest revenue and had the lowest cancellation rate (17.6%).",
    
    "2. Travel Agent bookings had the highest cancellation rate (27.9%), showing customers may change plans more often through this channel.",
    
    "3. Standard rooms had more cancellations (23.3%) compared to Deluxe rooms (16%).",
    
    "4. Booking volume was highest in April, which may indicate seasonal demand or holiday travel.",
    
    "5. Higher star hotels had slightly higher booking values, but cancellation rates remained similar across ratings.",
    
    "6. Coupon usage did not make a major difference in cancellation behavior."
]

for i, f in enumerate(findings):
    add_text_box(
        s,
        f,
        0.5,
        0.95 + i * 0.74,
        9,
        0.62,
        size=10.5,
        color=(60,60,60)
    )

# -- Slide 3: Chart 1 --
s = prs.slides.add_slide(blank)
bg = s.background; bg.fill.solid()
bg.fill.fore_color.rgb = rgb(244, 247, 252)
add_text_box(s, "Booking Status, Channel & Room Type", 0.3, 0.1, 9, 0.5,
             size=16, bold=True, color=(13,79,140))
s.shapes.add_picture("chart1_overview.png", Inches(0.2), Inches(0.7),
                     Inches(9.6), Inches(4.7))

# -- Slide 4: Chart 2 --
s = prs.slides.add_slide(blank)
bg = s.background; bg.fill.solid()
bg.fill.fore_color.rgb = rgb(244, 247, 252)
add_text_box(s, "Seasonal & Monthly Trends", 0.3, 0.1, 9, 0.5,
             size=16, bold=True, color=(13,79,140))
s.shapes.add_picture("chart2_trends.png", Inches(0.2), Inches(0.7),
                     Inches(9.6), Inches(4.7))

# -- Slide 5: Chart 3 --
s = prs.slides.add_slide(blank)
bg = s.background; bg.fill.solid()
bg.fill.fore_color.rgb = rgb(244, 247, 252)
add_text_box(s, "Revenue & Profitability", 0.3, 0.1, 9, 0.5,
             size=16, bold=True, color=(13,79,140))
s.shapes.add_picture("chart3_revenue.png", Inches(0.2), Inches(0.7),
                     Inches(9.6), Inches(4.7))

# -- Slide 6: Root Cause Analysis --
s = prs.slides.add_slide(blank)

bg = s.background
bg.fill.solid()
bg.fill.fore_color.rgb = rgb(244, 247, 252)

add_text_box(
    s,
    "Root Cause Analysis",
    0.4,
    0.2,
    9,
    0.55,
    size=20,
    bold=True,
    color=(13,79,140)
)

causes = [
    "1. Travel Agent bookings may have higher cancellations because customers often compare multiple hotel options before finalizing.",
    
    "2. Standard rooms may cancel more often because they are budget-friendly and customers may switch plans easily.",
    
    "3. Web bookings performed better because customers can directly compare prices, reviews, and availability.",
    
    "4. Higher booking volume in April may be influenced by holidays, vacations, or seasonal travel demand.",
    
    "5. Stay length and booking value may vary depending on travel season and customer preferences."
]

for i, c in enumerate(causes):
    add_text_box(
        s,
        c,
        0.5,
        0.9 + i * 0.7,
        9,
        0.6,
        size=10.5,
        color=(60,60,60)
    )

# -- Slide 7: Business Recommendations --
s = prs.slides.add_slide(blank)

bg = s.background
bg.fill.solid()
bg.fill.fore_color.rgb = rgb(244, 247, 252)

add_text_box(
    s,
    "Business Recommendations",
    0.4,
    0.2,
    9,
    0.55,
    size=20,
    bold=True,
    color=(13,79,140)
)

recs = [
    "1. Send booking reminders before check-in to help reduce cancellations.",
    
    "2. Offer small discounts or loyalty rewards for repeat customers.",
    
    "3. Focus more on Web and Mobile App bookings since they generate strong revenue.",
    
    "4. Provide seasonal offers during high-demand months to improve profitability.",
    
    "5. Improve Standard room packages with better deals to reduce cancellations."
]

for i, r in enumerate(recs):
    add_text_box(
        s,
        r,
        0.5,
        0.9 + i * 0.7,
        9,
        0.6,
        size=10.5,
        color=(60,60,60)
    )

# -- Slide 8: Final Summary --
s = prs.slides.add_slide(blank)

bg = s.background
bg.fill.solid()
bg.fill.fore_color.rgb = rgb(244, 247, 252)

add_text_box(
    s,
    "Final Summary",
    0.4,
    0.2,
    9,
    0.55,
    size=20,
    bold=True,
    color=(13,79,140)
)

summary = """
• Web channel performed best in terms of revenue and lower cancellations.

• Travel Agent bookings showed the highest cancellation rate.

• Booking activity was highest in April, showing seasonal demand.

• Standard rooms had higher cancellations compared to Deluxe rooms.

• Businesses can reduce cancellations through reminders, offers, and better customer engagement.
"""

add_text_box(
    s,
    summary,
    0.5,
    1.0,
    8.5,
    3.5,
    size=11,
    color=(60,60,60)
)

prs.save("Hotel_Bookings_Analysis.pptx")
print("Saved Hotel_Bookings_Analysis.pptx")
print("\nDone.")
